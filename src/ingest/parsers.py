"""Чтение входных файлов реального корпуса (разведка 02.07):
pdf (1163), docx (115), doc (18), docm (3), pptx (5), md/txt.
Для .doc нужен системный конвертер: catdoc или antiword или libreoffice
(apt install -y catdoc antiword). Без них .doc пропускается с предупреждением.
"""
from __future__ import annotations

import csv
import re
import shutil
import subprocess
from pathlib import Path

TEXT_EXT = {".md", ".txt"}
DOC_EXT = TEXT_EXT | {".pdf", ".docx", ".docm", ".doc", ".pptx"}


def read_document(path: str | Path) -> str:
    p = Path(path)
    ext = p.suffix.lower()
    if ext in TEXT_EXT:
        return p.read_text(encoding="utf-8", errors="ignore")
    if ext == ".pdf":
        return _read_pdf(p)
    if ext in {".docx", ".docm"}:  # docm — тот же zip-контейнер
        return _read_docx(p)
    if ext == ".doc":
        return _read_doc_legacy(p)
    if ext == ".pptx":
        return _read_pptx(p)
    raise ValueError(f"Неизвестный формат документа: {p}")


def _read_pdf(p: Path) -> str:
    import fitz  # pymupdf

    doc = fitz.open(p)
    pages = [f"[стр. {i + 1}]\n{page.get_text()}" for i, page in enumerate(doc)]
    return "\n\n".join(pages)


def _read_docx(p: Path) -> str:
    try:
        import docx

        d = docx.Document(str(p))
        parts = [par.text for par in d.paragraphs if par.text.strip()]
        for table in d.tables:
            for row in table.rows:
                parts.append(" | ".join(c.text.strip() for c in row.cells))
        return "\n".join(parts)
    except Exception:
        # .docm (macroEnabled) python-docx отклоняет по content-type — читаем zip напрямую
        return _read_docx_zip(p)


def _read_docx_zip(p: Path) -> str:
    import zipfile

    with zipfile.ZipFile(p) as z:
        xml = z.read("word/document.xml").decode("utf-8", errors="ignore")
    xml = xml.replace("</w:p>", "\n")
    return re.sub(r"<[^>]+>", "", xml)


def _read_doc_legacy(p: Path) -> str:
    """Старый бинарный .doc: catdoc -> antiword -> libreoffice."""
    for cmd in (["catdoc", "-w", str(p)], ["antiword", str(p)]):
        if shutil.which(cmd[0]):
            r = subprocess.run(cmd, capture_output=True, timeout=120)
            if r.returncode == 0 and len(r.stdout) > 200:
                return r.stdout.decode("utf-8", errors="ignore")
    if shutil.which("libreoffice"):
        out_dir = p.parent
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "txt:Text", "--outdir", str(out_dir), str(p)],
            capture_output=True, timeout=300,
        )
        txt = out_dir / (p.stem + ".txt")
        if txt.exists():
            text = txt.read_text(encoding="utf-8", errors="ignore")
            txt.unlink()
            return text
    raise RuntimeError(f".doc не сконвертирован (поставь catdoc/antiword): {p}")


def _read_pptx(p: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(p))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    texts.append(" | ".join(c.text.strip() for c in row.cells))
        if texts:
            parts.append(f"[слайд {i}]\n" + "\n".join(t for t in texts if t.strip()))
    return "\n\n".join(parts)


def clean_text(text: str) -> str:
    """Мусор Word-полей (координаты рамок), поля TOC/PAGEREF, лишние пробелы."""
    text = re.sub(r"(?:-?\d{4,}\s+){2,}-?\d{0,}", " ", text)  # "-813436 -396240 0 0 ..."
    text = re.sub(r"(?:TOC|PAGEREF|_Toc\d+)\s*\\?[a-z\\ \"';0-9]*", " ", text)
    text = re.sub(r"[ \t]+", " ", text)
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def read_catalog(path: str | Path) -> list[dict]:
    """Каталоги (эксперименты, материалы, оборудование, сотрудники) -> list[dict]."""
    p = Path(path)
    ext = p.suffix.lower()
    if ext == ".csv":
        with open(p, encoding="utf-8-sig") as f:
            return list(csv.DictReader(f))
    if ext in {".xlsx", ".xls"}:
        import pandas as pd

        return pd.read_excel(p).fillna("").to_dict(orient="records")
    raise ValueError(f"Неизвестный формат каталога: {p}")


def iter_documents(folder: str | Path):
    """Все документы из папки (md/txt/pdf/docx/docm/doc/pptx)."""
    folder = Path(folder)
    for p in sorted(folder.rglob("*")):
        if p.suffix.lower() in DOC_EXT and p.is_file():
            yield p
